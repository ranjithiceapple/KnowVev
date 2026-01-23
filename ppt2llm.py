from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
import argparse
import sys
import json
import re
from collections import Counter
from datetime import datetime
import pandas as pd
try:
    import tiktoken
except ImportError:
    tiktoken = None


def convert_pptx_to_markdown(pptx_path, output_path=None, include_notes=False):
    """
    Convert PPTX to Markdown format

    Args:
        pptx_path: Path to the input PPTX file
        output_path: Path for the output markdown file (optional)
        include_notes: Include speaker notes in output

    Returns:
        markdown_text: The converted markdown text
    """
    pptx_path = Path(pptx_path)

    # Validate input file
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    if not pptx_path.suffix.lower() in ['.pptx', '.ppt']:
        raise ValueError(f"File must be a PPTX: {pptx_path}")

    print(f"Converting PPTX to Markdown: {pptx_path.name}")
    print(f"{'='*70}")

    try:
        # Load presentation
        prs = Presentation(pptx_path)

        # Convert to markdown
        markdown_lines = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # Add slide separator
            markdown_lines.append(f"\n---\n")
            markdown_lines.append(f"<!-- Slide {slide_num} -->\n")

            # Extract slide content
            slide_content = extract_slide_content(slide, slide_num)

            if slide_content['title']:
                markdown_lines.append(f"# {slide_content['title']}\n")

            # Add text content
            for text in slide_content['text']:
                markdown_lines.append(text)

            # Add tables
            for table_md in slide_content['tables']:
                markdown_lines.append(table_md)
                markdown_lines.append('')

            # Add notes if requested
            if include_notes and slide_content['notes']:
                markdown_lines.append(f"\n**Notes:**\n")
                markdown_lines.append(f"> {slide_content['notes']}\n")

        markdown_text = '\n'.join(markdown_lines)

        # Determine output path
        if output_path is None:
            output_path = pptx_path.with_suffix('.md')
        else:
            output_path = Path(output_path)

        # Save markdown file
        print(f"Saving markdown to: {output_path.name}")
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown_text)

        # Calculate statistics
        lines = len(markdown_text.split('\n'))
        words = len(markdown_text.split())
        chars = len(markdown_text)

        print(f"{'='*70}")
        print("Conversion Statistics:")
        print(f"  Slides: {len(prs.slides)}")
        print(f"  Lines: {lines:,}")
        print(f"  Words: {words:,}")
        print(f"  Characters: {chars:,}")
        print(f"{'='*70}")
        print(f"SUCCESS: Markdown saved to {output_path}")

        return markdown_text

    except Exception as e:
        print(f"ERROR: Failed to convert PPTX")
        print(f"  {str(e)}")
        raise


def convert_pptx_to_text(pptx_path, output_path=None, include_notes=False):
    """
    Convert PPTX to plain text format

    Args:
        pptx_path: Path to the input PPTX file
        output_path: Path for the output text file (optional)
        include_notes: Include speaker notes in output

    Returns:
        text: The extracted plain text
    """
    pptx_path = Path(pptx_path)

    # Validate input file
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    print(f"Converting PPTX to Text: {pptx_path.name}")
    print(f"{'='*70}")

    try:
        # Load presentation
        prs = Presentation(pptx_path)

        # Extract text
        text_lines = []

        for slide_num, slide in enumerate(prs.slides, 1):
            text_lines.append(f"\n{'='*70}")
            text_lines.append(f"SLIDE {slide_num}")
            text_lines.append(f"{'='*70}\n")

            # Extract all text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_lines.append(shape.text.strip())
                    text_lines.append('')

            # Add notes if requested
            if include_notes and slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_lines.append(f"\nNotes: {notes_text}\n")

        text = '\n'.join(text_lines)

        # Determine output path
        if output_path is None:
            output_path = pptx_path.with_suffix('.txt')
        else:
            output_path = Path(output_path)

        # Save text file
        print(f"Saving text to: {output_path.name}")
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)

        # Statistics
        lines = len(text.split('\n'))
        words = len(text.split())

        print(f"{'='*70}")
        print("Extraction Statistics:")
        print(f"  Slides: {len(prs.slides)}")
        print(f"  Lines: {lines:,}")
        print(f"  Words: {words:,}")
        print(f"{'='*70}")
        print(f"SUCCESS: Text saved to {output_path}")

        return text

    except Exception as e:
        print(f"ERROR: Failed to convert PPTX")
        print(f"  {str(e)}")
        raise


def extract_slide_content(slide, slide_num):
    """
    Extract content from a single slide

    Returns:
        Dictionary with title, text, tables, and notes
    """
    content = {
        'slide_number': slide_num,
        'title': '',
        'text': [],
        'tables': [],
        'notes': ''
    }

    # Extract title
    if slide.shapes.title:
        content['title'] = slide.shapes.title.text.strip()

    # Extract text from shapes
    for shape in slide.shapes:
        # Skip title (already extracted)
        if shape == slide.shapes.title:
            continue

        # Extract text
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()

            # Check if it's a bullet list
            if hasattr(shape, "text_frame"):
                formatted_text = format_text_frame(shape.text_frame)
                if formatted_text:
                    content['text'].append(formatted_text)
            else:
                content['text'].append(text)

        # Extract tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_md = extract_table_from_shape(shape)
            if table_md:
                content['tables'].append(table_md)

    # Extract notes
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        if notes_text:
            content['notes'] = notes_text

    return content


def format_text_frame(text_frame):
    """
    Format text frame content as markdown (handles bullet points)
    """
    lines = []

    for paragraph in text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        # Get indentation level
        level = paragraph.level

        # Format as list item
        indent = '  ' * level
        lines.append(f"{indent}- {text}")

    return '\n'.join(lines) if lines else ''


def extract_table_from_shape(shape):
    """
    Extract table from a shape and convert to markdown
    """
    try:
        table = shape.table
        data = []

        # Extract table data
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            data.append(row_data)

        if not data or len(data) < 2:
            return None

        # Convert to DataFrame and markdown
        try:
            df = pd.DataFrame(data[1:], columns=data[0])
        except ValueError:
            df = pd.DataFrame(data)

        return df.to_markdown(index=False)

    except Exception as e:
        print(f"Warning: Could not extract table: {str(e)}")
        return None


def extract_pptx_tables_as_markdown(pptx_path):
    """
    Extract all tables from a PPTX file with slide context.
    Returns rich metadata for each table including slide information and context.

    Args:
        pptx_path: Path to PPTX file

    Returns:
        List of dictionaries containing:
            - markdown: Table in markdown format
            - slide_number: Slide number (1-indexed)
            - slide_title: Title of the slide
            - slide_context: All text from the slide
            - table_number: Table number on this slide
            - reference_name: Smart name combining slide and table info
            - full_content: Formatted string with source, context, and table
            - source_file: Original filename
            - headers: List of column headers
            - row_count: Number of data rows
            - dataframe: Original pandas DataFrame
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(pptx_path)
    file_name = pptx_path.name

    tables_with_context = []
    global_table_index = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        # 1. Gather Slide Context (Title + All Text)
        slide_text_elements = []
        slide_title = "Untitled Slide"
        
        # Extract slide title
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()
            if slide_title:
                slide_text_elements.append(f"Title: {slide_title}")
        
        # Gather all other text shapes for context
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                # Avoid duplicating the title
                if text and text != slide_title:
                    slide_text_elements.append(text)
        
        # Build full slide context
        slide_context_full = "\n".join(slide_text_elements) if slide_text_elements else ""
        
        # 2. Find and Extract Tables from this Slide
        table_num_on_slide = 0
        
        for shape in slide.shapes:
            # Check if shape is a table
            if shape.has_table or shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_num_on_slide += 1
                
                try:
                    table = shape.table
                    
                    # Extract table data
                    data = []
                    headers = None
                    
                    for row_idx, row in enumerate(table.rows):
                        row_data = []
                        for cell in row.cells:
                            # Extract text from cell's text frame
                            cell_text = cell.text.strip()
                            row_data.append(cell_text)
                        
                        # Skip completely empty rows
                        if not any(row_data):
                            continue
                        
                        # First non-empty row is headers
                        if row_idx == 0 or headers is None:
                            headers = row_data
                        else:
                            data.append(row_data)
                    
                    # Skip tables with no data
                    if not data:
                        continue
                    
                    # Convert to DataFrame → Markdown
                    try:
                        df = pd.DataFrame(data, columns=headers)
                        markdown_table = df.to_markdown(index=False)
                    except (ValueError, Exception) as e:
                        # Fallback if column mismatch
                        df = pd.DataFrame(data)
                        markdown_table = df.to_markdown(index=False)
                        headers = [f"Column_{i}" for i in range(len(data[0]))] if data else []
                    
                    # Create smart reference name
                    if table_num_on_slide == 1:
                        reference_name = f"Table on Slide {slide_num}: {slide_title}"
                    else:
                        reference_name = f"Table {table_num_on_slide} on Slide {slide_num}: {slide_title}"
                    
                    # Build rich content string
                    rich_content = f"""SOURCE FILE: {file_name}
TYPE: PPTX
SLIDE: {slide_num}
TITLE: {slide_title}

--- SLIDE CONTEXT ---
{slide_context_full}

--- TABLE DATA ---
{markdown_table}"""
                    
                    # Store table with comprehensive metadata
                    tables_with_context.append({
                        'markdown': markdown_table,
                        'slide_number': slide_num,
                        'slide_title': slide_title,
                        'slide_context': slide_context_full,
                        'table_number': table_num_on_slide,
                        'global_table_index': global_table_index,
                        'reference_name': reference_name,
                        'full_content': rich_content,
                        'source_file': file_name,
                        'headers': headers if headers else [],
                        'row_count': len(data),
                        'dataframe': df
                    })
                    
                    global_table_index += 1
                    
                except Exception as e:
                    print(f"Warning: Could not extract table from slide {slide_num}, table {table_num_on_slide}: {str(e)}")
                    continue

    return tables_with_context


def extract_pptx_with_structure(pptx_path):
    """
    Extract PPTX content preserving the order of slides, text, tables, and notes.
    Returns a unified list of content blocks with metadata for RAG.

    Args:
        pptx_path: Path to PPTX file

    Returns:
        List of content blocks, each with:
        {
            'type': 'title'|'text'|'table'|'notes'|'list',
            'content': str,
            'slide': int,
            'position': int (order on slide),
            'metadata': dict
        }
    """
    content_blocks = []
    pptx_path = Path(pptx_path)

    try:
        prs = Presentation(pptx_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            position = 0

            # Extract slide title
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    content_blocks.append({
                        'type': 'title',
                        'content': title_text,
                        'slide': slide_num,
                        'position': position,
                        'metadata': {
                            'pptx_file': pptx_path.name,
                            'slide': slide_num,
                            'block_type': 'title'
                        }
                    })
                    position += 1

            # Process all shapes in order
            for shape in slide.shapes:
                # Skip title (already processed)
                if shape == slide.shapes.title:
                    continue

                # Extract tables
                if shape.has_table:
                    table = shape.table
                    rows_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        rows_data.append(row_data)

                    if rows_data and len(rows_data) > 1:
                        try:
                            df = pd.DataFrame(rows_data[1:], columns=rows_data[0])
                            table_md = df.to_markdown(index=False)

                            content_blocks.append({
                                'type': 'table',
                                'content': table_md,
                                'slide': slide_num,
                                'position': position,
                                'metadata': {
                                    'pptx_file': pptx_path.name,
                                    'slide': slide_num,
                                    'block_type': 'table',
                                    'rows': len(rows_data),
                                    'columns': len(rows_data[0]) if rows_data else 0
                                }
                            })
                            position += 1
                        except Exception as e:
                            print(f"Warning: Could not convert table on slide {slide_num}: {e}")

                # Extract text from text frames
                elif shape.has_text_frame:
                    text_content = []
                    is_list = False

                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            # Check if it's a list
                            if paragraph.level > 0 or para_text.startswith(('•', '-', '*', '○')):
                                is_list = True
                                indent = "  " * paragraph.level
                                text_content.append(f"{indent}- {para_text}")
                            else:
                                text_content.append(para_text)

                    if text_content:
                        block_text = '\n'.join(text_content)
                        block_type = 'list' if is_list else 'text'

                        content_blocks.append({
                            'type': block_type,
                            'content': block_text,
                            'slide': slide_num,
                            'position': position,
                            'metadata': {
                                'pptx_file': pptx_path.name,
                                'slide': slide_num,
                                'block_type': block_type
                            }
                        })
                        position += 1

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                notes_text = notes_frame.text.strip() if notes_frame else ""

                if notes_text:
                    content_blocks.append({
                        'type': 'notes',
                        'content': f"Speaker Notes:\n{notes_text}",
                        'slide': slide_num,
                        'position': position,
                        'metadata': {
                            'pptx_file': pptx_path.name,
                            'slide': slide_num,
                            'block_type': 'notes'
                        }
                    })
                    position += 1

    except Exception as e:
        print(f"Error extracting PPTX with structure: {str(e)}")
        return []

    return content_blocks


def extract_pptx_metadata(pptx_path):
    """
    Extract comprehensive metadata from PPTX file

    Returns:
        Dictionary with presentation metadata
    """
    pptx_path = Path(pptx_path)
    prs = Presentation(pptx_path)

    # Core properties
    core_props = prs.core_properties

    metadata = {
        'filename': pptx_path.name,
        'file_size_bytes': pptx_path.stat().st_size,
        'file_size_mb': round(pptx_path.stat().st_size / (1024 * 1024), 2),
        'presentation_properties': {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
            'keywords': core_props.keywords or '',
            'comments': core_props.comments or '',
            'last_modified_by': core_props.last_modified_by or '',
            'created': str(core_props.created) if core_props.created else '',
            'modified': str(core_props.modified) if core_props.modified else '',
            'category': core_props.category or '',
        },
        'structure': {
            'slides': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
        },
        'slides_info': [],
        'total_shapes': 0,
        'total_tables': 0,
        'total_images': 0,
        'total_charts': 0,
        'slides_with_notes': 0
    }

    # Extract per-slide metadata
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': slide_num,
            'shapes': len(slide.shapes),
            'has_title': slide.shapes.title is not None,
            'title': slide.shapes.title.text if slide.shapes.title else '',
            'tables': 0,
            'images': 0,
            'charts': 0,
            'text_boxes': 0,
            'has_notes': slide.has_notes_slide
        }

        # Count shape types
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                slide_info['tables'] += 1
                metadata['total_tables'] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_info['images'] += 1
                metadata['total_images'] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                slide_info['charts'] += 1
                metadata['total_charts'] += 1
            elif hasattr(shape, "text"):
                slide_info['text_boxes'] += 1

        if slide.has_notes_slide:
            metadata['slides_with_notes'] += 1

        metadata['total_shapes'] += len(slide.shapes)
        metadata['slides_info'].append(slide_info)

    return metadata


def extract_slide_titles(markdown_text):
    """
    Extract slide titles from markdown text
    """
    titles = []
    lines = markdown_text.split('\n')

    current_slide = 0
    for line in lines:
        # Detect slide markers
        if '<!-- Slide' in line:
            match = re.search(r'Slide (\d+)', line)
            if match:
                current_slide = int(match.group(1))

        # Extract H1 headings (slide titles)
        if line.startswith('# '):
            title = line[2:].strip()
            titles.append({
                'slide_number': current_slide,
                'title': title
            })

    return titles


def analyze_presentation_structure(markdown_text):
    """
    Analyze presentation structure from markdown
    """
    lines = markdown_text.split('\n')

    structure = {
        'total_lines': len(lines),
        'slides': markdown_text.count('<!-- Slide'),
        'titles': len([l for l in lines if l.startswith('# ')]),
        'lists': {
            'unordered': len(re.findall(r'^\s*[-*+]\s+', markdown_text, re.MULTILINE)),
        },
        'tables': len(re.findall(r'\|[^\n]+\|', markdown_text)),
        'notes': markdown_text.count('**Notes:**'),
    }

    # Word statistics
    words = markdown_text.split()
    structure['word_count'] = len(words)
    structure['char_count'] = len(markdown_text)
    structure['avg_words_per_slide'] = round(len(words) / structure['slides'], 2) if structure['slides'] > 0 else 0

    return structure


def generate_metadata_report(pptx_path, markdown_text, output_path=None):
    """
    Generate comprehensive metadata report for PPTX

    Args:
        pptx_path: Path to PPTX file
        markdown_text: Converted markdown text
        output_path: Optional path for metadata JSON file

    Returns:
        Complete metadata dictionary
    """
    print(f"\n{'='*70}")
    print("EXTRACTING METADATA AND PRESENTATION STRUCTURE")
    print(f"{'='*70}\n")

    # Extract all metadata
    print("Extracting PPTX metadata...")
    pptx_metadata = extract_pptx_metadata(pptx_path)

    print("Analyzing slide titles...")
    titles = extract_slide_titles(markdown_text)

    print("Analyzing presentation structure...")
    structure = analyze_presentation_structure(markdown_text)

    # Compile complete report
    report = {
        'extraction_date': datetime.now().isoformat(),
        'source_file': str(pptx_path),
        'pptx_metadata': pptx_metadata,
        'slide_titles': {
            'count': len(titles),
            'list': titles
        },
        'structure': structure
    }

    # Save to JSON if output path provided
    if output_path:
        output_path = Path(output_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
        print(f"\nMetadata saved to: {output_path}")

    return report


def print_metadata_summary(report):
    """
    Print a human-readable summary of the metadata report
    """
    print(f"\n{'='*70}")
    print("METADATA SUMMARY")
    print(f"{'='*70}\n")

    # Presentation Info
    pptx = report['pptx_metadata']
    print(f"Presentation Information:")
    print(f"  File: {pptx['filename']}")
    print(f"  Size: {pptx['file_size_mb']} MB")
    print(f"  Slides: {pptx['structure']['slides']}")
    print(f"  Total Shapes: {pptx['total_shapes']}")

    # Presentation Properties
    props = pptx['presentation_properties']
    if any(props.values()):
        print(f"\nPresentation Properties:")
        if props['title']:
            print(f"  Title: {props['title']}")
        if props['author']:
            print(f"  Author: {props['author']}")
        if props['subject']:
            print(f"  Subject: {props['subject']}")
        if props['created']:
            print(f"  Created: {props['created']}")
        if props['modified']:
            print(f"  Modified: {props['modified']}")

    # Content Statistics
    print(f"\nContent Statistics:")
    print(f"  Tables: {pptx['total_tables']}")
    print(f"  Images: {pptx['total_images']}")
    print(f"  Charts: {pptx['total_charts']}")
    print(f"  Slides with Notes: {pptx['slides_with_notes']}")

    # Titles
    titles = report['slide_titles']
    print(f"\nSlide Titles: {titles['count']} slides with titles")

    # Structure
    struct = report['structure']
    print(f"\nPresentation Structure:")
    print(f"  Total Slides: {struct['slides']}")
    print(f"  Words: {struct['word_count']:,}")
    print(f"  Avg Words/Slide: {struct['avg_words_per_slide']}")
    print(f"  Lists: {struct['lists']['unordered']} items")
    print(f"  Tables: {struct['tables']}")

    # Slide Preview
    if titles['list']:
        print(f"\nFirst 5 Slide Titles:")
        for slide in titles['list'][:5]:
            print(f"  Slide {slide['slide_number']}: {slide['title']}")

    print(f"\n{'='*70}")


def batch_convert_pptx(pptx_directory, output_directory=None, output_format='md', include_notes=False):
    """
    Convert all PPTX files in a directory

    Args:
        pptx_directory: Directory containing PPTX files
        output_directory: Directory for output files (optional)
        output_format: 'md' for markdown or 'txt' for plain text
        include_notes: Include speaker notes
    """
    pptx_dir = Path(pptx_directory)

    if not pptx_dir.exists() or not pptx_dir.is_dir():
        raise ValueError(f"Invalid directory: {pptx_directory}")

    # Find all PPTX files
    pptx_files = list(pptx_dir.glob('*.pptx'))

    if not pptx_files:
        print(f"No PPTX files found in: {pptx_directory}")
        return

    print(f"\nFound {len(pptx_files)} PPTX files")
    print(f"{'='*70}\n")

    # Setup output directory
    if output_directory:
        output_dir = Path(output_directory)
        output_dir.mkdir(parents=True, exist_ok=True)
    else:
        output_dir = pptx_dir

    # Convert each PPTX
    success_count = 0
    failed_files = []

    for i, pptx_file in enumerate(pptx_files, 1):
        print(f"\n[{i}/{len(pptx_files)}] Converting: {pptx_file.name}")

        try:
            if output_format == 'md':
                output_file = output_dir / pptx_file.with_suffix('.md').name
                convert_pptx_to_markdown(pptx_file, output_file, include_notes)
            else:
                output_file = output_dir / pptx_file.with_suffix('.txt').name
                convert_pptx_to_text(pptx_file, output_file, include_notes)
            success_count += 1
        except Exception as e:
            print(f"  FAILED: {str(e)}")
            failed_files.append(pptx_file.name)

    # Summary
    print(f"\n{'='*70}")
    print(f"BATCH CONVERSION COMPLETE")
    print(f"{'='*70}")
    print(f"  Successful: {success_count}/{len(pptx_files)}")

    if failed_files:
        print(f"  Failed files:")
        for fname in failed_files:
            print(f"    - {fname}")


def main():
    """
    Command-line interface for PPTX conversion
    """
    parser = argparse.ArgumentParser(
        description='Convert PPTX files to Markdown or Text',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
Examples:
  # Convert to markdown
  python ppt2llm.py presentation.pptx

  # Convert to plain text
  python ppt2llm.py presentation.pptx --format txt

  # Include speaker notes
  python ppt2llm.py presentation.pptx --notes

  # Extract metadata only
  python ppt2llm.py presentation.pptx --metadata-only

  # Batch convert all PPTX files
  python ppt2llm.py --batch ./presentations/ -o ./output/
        '''
    )

    parser.add_argument(
        'input',
        nargs='?',
        help='Input PPTX file or directory (for batch mode)'
    )

    parser.add_argument(
        '-o', '--output',
        help='Output file or directory (for batch mode)'
    )

    parser.add_argument(
        '--format',
        choices=['md', 'txt'],
        default='md',
        help='Output format: md (markdown) or txt (plain text)'
    )

    parser.add_argument(
        '--notes',
        action='store_true',
        help='Include speaker notes in output'
    )

    parser.add_argument(
        '--batch',
        action='store_true',
        help='Batch convert all PPTX files in input directory'
    )

    parser.add_argument(
        '--metadata',
        action='store_true',
        help='Extract and save metadata'
    )

    parser.add_argument(
        '--metadata-only',
        action='store_true',
        help='Extract only metadata without conversion'
    )

    args = parser.parse_args()

    # Show help if no input provided
    if not args.input:
        parser.print_help()
        sys.exit(1)

    try:
        if args.metadata_only:
            # Extract metadata only
            print("Extracting metadata only...")
            pptx_path = Path(args.input)

            # Convert to get markdown for analysis
            markdown_text = convert_pptx_to_markdown(pptx_path, include_notes=args.notes)

            # Generate metadata report
            metadata_output = args.output or pptx_path.with_suffix('.metadata.json')
            report = generate_metadata_report(pptx_path, markdown_text, metadata_output)
            print_metadata_summary(report)

        elif args.batch:
            # Batch conversion mode
            batch_convert_pptx(
                args.input,
                args.output,
                args.format,
                args.notes
            )
        else:
            # Single file conversion
            if args.format == 'md':
                markdown_text = convert_pptx_to_markdown(args.input, args.output, args.notes)
            else:
                markdown_text = convert_pptx_to_text(args.input, args.output, args.notes)

            # Extract metadata if requested
            if args.metadata:
                pptx_path = Path(args.input)
                metadata_output = pptx_path.with_suffix('.metadata.json')
                report = generate_metadata_report(pptx_path, markdown_text, metadata_output)
                print_metadata_summary(report)

    except KeyboardInterrupt:
        print("\n\nConversion cancelled by user")
        sys.exit(1)
    except Exception as e:
        print(f"\nERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
